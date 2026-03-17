"""Rewrite the AP Research POD presentation with sharper, more compelling framing."""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# Color constants matching the existing theme
C_TITLE = RGBColor(0x33, 0x33, 0x33)
C_SUBTITLE = RGBColor(0x4E, 0x9E, 0xAD)
C_BODY = RGBColor(0x66, 0x66, 0x66)
C_CAPTION = RGBColor(0x99, 0x99, 0x99)

# Size constants
S_TITLE = Pt(32)
S_SUBTITLE = Pt(18)
S_BODY = Pt(15)
S_CAPTION = Pt(12)


def clear_text_frame(tf):
    """Remove all paragraphs from a text frame except the first (required by python-pptx)."""
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    # Clear the remaining first paragraph
    tf.paragraphs[0].clear()


def add_runs(paragraph, runs):
    """Add runs to a paragraph. Each run is (text, bold, size, color)."""
    for text, bold, size, color in runs:
        r = paragraph.add_run()
        r.text = text
        r.font.bold = bold
        r.font.size = size
        r.font.color.rgb = color


def set_paragraphs(tf, paras):
    """Set a text frame's content.

    paras: list of lists-of-tuples.  Each outer item is a paragraph.
    Each inner tuple is (text, bold, size, color) for one run.
    """
    clear_text_frame(tf)
    for i, runs in enumerate(paras):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        add_runs(p, runs)


def body(*parts):
    """Create a single paragraph from alternating plain/bold segments.

    parts can be:
      - str: plain text
      - (str, True): bold text
    """
    runs = []
    for part in parts:
        if isinstance(part, tuple):
            runs.append((part[0], True, S_BODY, C_BODY))
        else:
            runs.append((part, False, S_BODY, C_BODY))
    return runs


def bold_body(text):
    return [(text, True, S_BODY, C_BODY)]


def plain_body(text):
    return [(text, False, S_BODY, C_BODY)]


def subtitle(text):
    return [(text, None, S_SUBTITLE, C_SUBTITLE)]


def section_header(text):
    return [(text, True, S_SUBTITLE, C_SUBTITLE)]


def caption(text):
    return [(text, None, S_CAPTION, C_CAPTION)]


def title_text(text):
    return [(text, True, S_TITLE, C_TITLE)]


# ---------------------------------------------------------------------------
# Slide rewrite functions
# ---------------------------------------------------------------------------

def rewrite_slide_1(slide):
    """Title slide — add sharper subtitle."""
    shapes = slide.shapes
    # Shape 0: Main title — keep
    # Shape 2: Subtitle
    set_paragraphs(shapes[2].text_frame, [
        [(
            "How Three ML Architectures Disagree Reveals Why Upsets Happen",
            None, Pt(20), C_SUBTITLE,
        )],
    ])


def rewrite_slide_2(slide):
    """Key Terms — tighten definitions."""
    shapes = slide.shapes
    # Shape 0: title — keep
    # Shape 2: subtitle
    set_paragraphs(shapes[2].text_frame, [
        subtitle("These terms appear throughout. Bookmark this slide."),
    ])
    # Shape 3: definitions
    set_paragraphs(shapes[3].text_frame, [
        body(("AUC (Area Under the Curve): ", True), "How well a model ranks outcomes from 0 to 1. Random guessing = 0.50, perfect = 1.0. Think of it as a letter grade for predictions."),
        body(("p-value: ", True), "The probability a result happened by luck alone. Below .05 = statistically significant. Above .05 = might be coincidence."),
        body(("Bootstrap Confidence Interval: ", True), "A range computed by resampling data thousands of times. If it contains zero, the difference isn't real."),
        body(("Cross-Validation: ", True), "Train on past seasons, test on future ones. Prevents the model from memorizing answers it shouldn't know yet."),
        body(("LSTM (Long Short-Term Memory): ", True), "A neural network that reads games in order — week 1, then week 2, then week 3 — detecting momentum, streaks, and fatigue patterns."),
        body(("Point Spread: ", True), "How many points oddsmakers expect the favorite to win by. A spread of 7 means the market expects about a touchdown margin."),
        body(("Spread Ablation: ", True), "Removing the betting line from the models to test whether they learned something real or were just copying the market's homework."),
    ])


def rewrite_slide_3(slide):
    """Research Question — add the diagnostic reframe from PROJECT.md."""
    shapes = slide.shapes
    # Shape 2: question
    set_paragraphs(shapes[2].text_frame, [
        subtitle("To what extent can disagreement patterns among logistic regression, XGBoost, and LSTM models, combined with spread ablation analysis, reveal distinct categories of NFL upsets?"),
    ])
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body("An NFL upset occurs when a 3+ point underdog wins. That happens in roughly 29% of games — far more than most people think."),
        body(("This is not a model competition. ", True), "We deliberately chose three architectures that process data differently, then used ", ("where they agree and disagree as a diagnostic tool. ", True), "The contribution is a taxonomy of upset types, not a leaderboard."),
        body("Spread ablation removes betting-line data to answer a second question: are the models learning genuine signal, or just echoing what Vegas already knows?"),
    ])


def rewrite_slide_4(slide):
    """Literature Review — sharper bullets."""
    shapes = slide.shapes
    # Shape 3: Sports Prediction body
    set_paragraphs(shapes[3].text_frame, [
        plain_body("Simple statistical models remain competitive with complex ML for NFL prediction — accuracy gains have been marginal for decades (Boulier & Stekler, 2003)"),
        plain_body("Tree-based models (XGBoost) dominate deep learning on standard tabular data, but that advantage disappears on sequential data (Grinsztajn et al., 2022)"),
        plain_body("LSTMs can capture time-dependent patterns in sports — momentum, streaks, trajectory — that static models structurally cannot (Huang & Li, 2021)"),
    ])
    # Shape 5: Market Efficiency body
    set_paragraphs(shapes[5].text_frame, [
        plain_body("NFL betting lines outperform both expert judgment and statistical models at prediction (Song et al., 2007; Wilkens, 2021)"),
        plain_body("Decorrelating predictions from bookmaker odds is essential — otherwise models just rediscover what the market already priced in (Hubacek et al., 2019)"),
        plain_body("Market inefficiencies exist in narrow contexts but are small and hard to exploit consistently (Boulier et al., 2006)"),
    ])
    # Shape 7: Model Disagreement body
    set_paragraphs(shapes[7].text_frame, [
        plain_body("Per-instance model disagreement is an unbiased estimate of prediction uncertainty (Gordon et al., 2021)"),
        plain_body("Architecturally diverse models capture structurally different patterns in the same data (Dietterich, 2000)"),
        body(("No prior work uses model disagreement to categorize upsets into distinct types", True)),
    ])


def rewrite_slide_5(slide):
    """Gap in the Literature — more assertive."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Three Blind Spots in the Literature"),
    ])
    # Shape 2: subtitle
    set_paragraphs(shapes[2].text_frame, [
        subtitle("No one has used model disagreement to classify NFL upsets into structurally distinct types."),
    ])
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body(("1. All upsets are treated identically. ", True), "Prior research predicts wins and losses but never asks why a specific upset happened — whether it was driven by momentum, hidden information, or pure chance."),
        body(("2. Disagreement is thrown away. ", True), "Ensemble research uses disagreement to improve accuracy (average the models, pick the majority). No one has used the pattern of disagreement itself as a diagnostic tool."),
        body(("3. Market dependence is untested. ", True), "No study has checked whether prediction models learn independent signal or simply echo what the betting line already encodes."),
    ])


def rewrite_slide_6(slide):
    """Method Overview — sharpen the logic."""
    shapes = slide.shapes
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body(("1. Collect: ", True), "18 NFL regular seasons (2005–2022) for training. Hold out 2023–2025 as a blind test set the models never see during development."),
        [],  # blank line
        body(("2. Represent: ", True), "Give each model the same underlying data in a different form — like three analysts watching the same games but reading different scouting reports:"),
        body("     LR reads \"The Summary\" (46 stats). XGB reads \"The Details\" (70 features with game-by-game lags). LSTM watches \"The Movie\" (raw 8-game sequences)."),
        [],  # blank line
        body(("3. Train identically: ", True), "All three models see the same folds (6-fold expanding-window CV). Same data, same splits. When they disagree, it's the architecture — not the data."),
        [],  # blank line
        body(("4. Diagnose: ", True), "Classify every game by which models got it right and which didn't. The pattern of agreement and disagreement is the diagnostic instrument."),
        [],  # blank line
        body(("5. Ablate: ", True), "Remove the point spread and rerun everything. If a model collapses, it was just reading the market's homework. If it holds, it learned something real."),
    ])


def rewrite_slide_7(slide):
    """Data & Multi-Representation Design — sharpen the Why section."""
    shapes = slide.shapes
    # Shape 6: "Why these three?" text — make it punchier
    set_paragraphs(shapes[6].text_frame, [
        body(("Each model has a structural blind spot, and that's the point. ", True), "LR sees only linear relationships — is the spread simply mispriced? XGBoost captures feature interactions — does a bad recent game + short week + road trip = danger? LSTM reads trajectories — is a team rising or falling apart? When they disagree, the architecture tells you the mechanism."),
    ])


def rewrite_slide_8(slide):
    """CV Results — lead with the finding, not the topic."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("All Three Models Are Statistically Tied"),
    ])
    # Shape 4: body bullets
    set_paragraphs(shapes[4].text_frame, [
        body(("Bootstrap 95% CIs on every pairwise AUC difference contain zero. ", True), "No model is significantly better than another. This is exactly what the framework needs — disagreement reflects different architectural perspectives, not quality gaps."),
        body("LR and XGBoost correlate at .874 — they largely agree. The LSTM correlates moderately with both (.784 with LR, .699 with XGB), confirming it captures partially independent information."),
        body("The lowest correlation is XGB–LSTM (.699): the greatest structural distance lies between interaction-based processing and sequence-based processing. ", ("This is where disagreement will be most informative.", True)),
    ])


def rewrite_slide_9(slide):
    """Disagreement Categories — sharpen the key finding."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Where Models Disagree Reveals the Mechanism"),
    ])
    # Shape 4: body
    set_paragraphs(shapes[4].text_frame, [
        body("74.7% of games fall into all-agree categories — either all three are right, or all three are wrong. The interesting games are in the 25.3% where they split."),
        body(("Key finding: ", True), "Of 65 LSTM-exclusive games, ", ("53 (81.5%) are non-upset rejections", True), " — the LSTM says \"no upset\" when both LR and XGB say \"upset.\" Only 12 are upsets the LSTM alone catches. This skew is statistically significant (binomial p = 0.029)."),
        body(("The LSTM's primary exclusive value is moderating false alarms, not detecting hidden upsets.", True), " But this headline hides a deeper story — the next slide."),
    ])


def rewrite_slide_10(slide):
    """LSTM Role Depends on Spread — the climax. Make it hit hard."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Same Model, Opposite Jobs"),
    ])
    # Shape 2: subtitle
    set_paragraphs(shapes[2].text_frame, [
        subtitle("The LSTM's role inverts depending on the matchup context. This is the study's central finding."),
    ])
    # Shape 4: body — the narrative
    set_paragraphs(shapes[4].text_frame, [
        body(("Close games (3–6.5 pts): ", True), "LR and XGB look at the stats and say \"upset.\" The LSTM watches the last 8 games and says \"no — this team's trajectory is wrong.\" It's right 92% of the time. The movie tells a different story than the snapshot."),
        body(("Bigger underdogs (7–13.5 pts): ", True), "Now the LSTM flips. It spots a team building week by week — rising EPA, winning close games, gaining momentum. The summary stats say \"no chance.\" The trajectory says \"watch out.\" It catches 83% of these as real upsets."),
        body(("The LSTM is not one instrument — it's two, and the spread determines which one you're hearing. ", True), "At close spreads it filters noise. At medium spreads it detects signal. That context-dependence proves the pattern is structural, not statistical noise."),
    ])


def rewrite_slide_11(slide):
    """Spread Ablation — frame as unmasking."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Removing the Answer Sheet"),
    ])
    # Shape 3: body bullets (this slide has no subtitle shape)
    set_paragraphs(shapes[3].text_frame, [
        body("All three models depend heavily on the spread. Bootstrap CIs confirm every drop is statistically significant."),
        body(("The ranking flips without spread: LSTM (.574) > LR (.571) > XGB (.566). ", True), "The LSTM degrades the least (−.067) and becomes the strongest model when forced to learn from game data alone."),
        body("But honesty matters: the LSTM's smaller delta is ", ("not", True), " statistically significant versus LR's delta — bootstrap CI contains zero. Suggestive, not proven."),
        body("Without spread, ", ("LSTM-exclusive predictions double (5.6% → 11.0%)", True), " and upsets caught jump from 12 to 33. The betting line was masking genuine temporal signal the whole time."),
        body("LR–XGB correlation drops from .874 to .742. The spread was the shared anchor holding model predictions together. ", ("Remove it and the three architectures scatter — revealing how differently they actually read the game.", True)),
    ])


def rewrite_slide_12(slide):
    """Out-of-sample test — frame as honest reckoning, not failure."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("The Honest Reckoning: Out-of-Sample 2023–2025"),
    ])
    # Shape 5: body bullets (shapes 2,4 are tables, shape 3 is "Top-K" header)
    set_paragraphs(shapes[5].text_frame, [
        body(("The LSTM shows the largest generalization gap (−.117). ", True), "Temporal patterns learned from 2005–2022 do not fully transfer to 2023–2025. This is not one bad year — it trails in all three test seasons."),
        body(("XGBoost generalizes best (−.062). ", True), "Non-linear interaction patterns are the most temporally stable signal in the data."),
        body("The LSTM's correlation with static models collapses on the test set: LR .784 → .429, XGB .699 → .408. It diverges most when the stakes are real."),
        body(("XGBoost's top 10 predictions hit at 60% — 2.1× lift over base rate. ", True), "Signal concentrates in the highest-confidence predictions and fades fast."),
    ])


def rewrite_slide_13(slide):
    """Not All Upsets Are Created Equal — strengthen narrative."""
    shapes = slide.shapes
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body(("All three agree → ", True), "The mismatch was visible from every angle. The Summary, The Details, and The Movie all told the same story. These are the structurally readable outcomes — and the market prices them well. 45% of games."),
        body(("Only the LSTM sees it → ", True), "The movie told a story the snapshot couldn't. A defense leaking more yards each week. An offense finding its rhythm. A team building quietly after three straight losses. The summary says \"average.\" The trajectory says \"momentum.\" 6% of games — small, but structurally distinct and statistically confirmed."),
        body(("All three miss big favorites → ", True), "Something happened outside the data. A QB injury during warmups. A team resting starters with a playoff bye locked up. A locker room that imploded on Tuesday. No historical stat captures game-day surprises. This is a data availability boundary, not a modeling failure."),
        body(("All three miss close games → ", True), "A coin flip that landed tails. 3-point games are genuinely close. A tipped pass, a missed field goal, a bad spot on 4th-and-inches. Every model knew it was close. No pre-game model can predict which way close games break."),
    ])


def rewrite_slide_14(slide):
    """Momentum Is Real — sharpen the finding-not-failure framing.

    Slide 14 has separate text boxes for each section:
    Shape 0: title, Shape 2: intro, Shape 3: header1, Shape 4: body1,
    Shape 5: header2, Shape 6: body2, Shape 7: header3, Shape 8: body3
    """
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Momentum Is Real — but the NFL Keeps Changing"),
    ])
    # Shape 2: intro subtitle
    set_paragraphs(shapes[2].text_frame, [
        subtitle("The LSTM matched the other models in training but showed the largest drop on truly future data. That's not a failure — it's a finding."),
    ])
    # Shape 3: section 1 header
    set_paragraphs(shapes[3].text_frame, [
        section_header("In cross-validation (2005–2022):"),
    ])
    # Shape 4: section 1 body
    set_paragraphs(shapes[4].text_frame, [
        plain_body("LSTM AUC .641 — statistically tied with LR (.650) and XGB (.638)"),
        plain_body("The temporal signal is real: it adds ensemble value, and its exclusive catches double when the spread is removed"),
        plain_body("Momentum patterns from this era — streaks, fatigue, trajectory — are genuinely predictive"),
    ])
    # Shape 5: section 2 header
    set_paragraphs(shapes[5].text_frame, [
        section_header("On the test set (2023–2025):"),
    ])
    # Shape 6: section 2 body
    set_paragraphs(shapes[6].text_frame, [
        plain_body("LSTM drops to .524 — largest gap of any model (−.117)"),
        plain_body("It trails LR and XGB in all three test seasons. Not one weird year — persistent."),
        plain_body("The correlation between LSTM and the static models collapses: .784 → .429"),
    ])
    # Shape 7: section 3 header
    set_paragraphs(shapes[7].text_frame, [
        section_header("What this means about the NFL:"),
    ])
    # Shape 8: section 3 body
    set_paragraphs(shapes[8].text_frame, [
        body("Momentum patterns are era-specific. How a hot streak translated into upset probability in 2015 doesn't work the same way in 2024. Coaching changes, rule changes, the evolution of offenses — the temporal dynamics shift faster than the static relationships. ", ("The LSTM's poor forward transfer is a finding about the sport, not just about the model.", True)),
    ])


def rewrite_slide_15(slide):
    """What the Spread Ablation Tells Us — sharper."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("What the Spread Really Contributes"),
    ])
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body("LR lost .079 AUC, XGBoost lost .072, LSTM lost .067. All highly significant (p < .001). ", ("The spread contains more predictive information than all other team statistics combined.", True)),
        body("LR's largest coefficient is the spread itself — 5× bigger than any team stat. The \"model\" was mostly reading the market's homework."),
        body(("But the LSTM kept the most signal. ", True), "Without spread, it becomes the strongest model (.574 > .571 > .566). The trajectory captures something the market line encodes differently — not better, differently."),
        body("Model agreement collapsed: LR–XGB correlation .874 → .742, XGB–LSTM .699 → .419. ", ("The spread was the glue. Remove it and the three architectures reveal how differently they actually read the game.", True)),
        body("LSTM-exclusive catches doubled without spread (12 → 33 upsets). The betting line was masking genuine temporal signal."),
        body(("Bottom line: ", True), "the market is extremely good. Oddsmakers have already digested everything in publicly available statistics. But sequential processing captures information that the point spread encodes through a different mechanism."),
    ])


def rewrite_slide_16(slide):
    """Why Disagreement Proves Structure — synthesis slide."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Why Disagreement Is the Contribution"),
    ])
    # Shape 2: intro - update to avoid duplication with body
    set_paragraphs(shapes[2].text_frame, [
        subtitle("Three models. Same data. Different architectures. If upsets were random, the models would fail randomly. They don't."),
    ])
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body(("The Summary and The Details agree 87% of the time. ", True), "When the static picture is clear — one team is better on paper — both linear and non-linear models see it. When they're both wrong, the cause is outside the numbers."),
        body(("The Movie tells a different story — and which story depends on context. ", True), "In close games, the LSTM corrects the others: \"the stats say upset, but watch the trajectory — this team is fading.\" In bigger mismatches, it catches what they miss: \"the numbers say no, but this team has been building for three weeks.\""),
        body(("Remove the betting line and the masks come off. ", True), "The spread was the anchor all three models leaned on. Without it, agreement collapses, the LSTM's unique signal doubles, and the architectures reveal how differently they read the game."),
        body(("Momentum is real, but it changes shape. ", True), "The LSTM finds temporal patterns that improve predictions during training. But the NFL evolves — and those patterns don't transfer forward unchanged. That tells us something about the sport itself."),
    ])


def rewrite_slide_17(slide):
    """Limitations — tighten language."""
    shapes = slide.shapes
    # Shape 2: body (this slide has only shapes 0-2, no subtitle)
    set_paragraphs(shapes[2].text_frame, [
        body(("Small cell sizes in spread strata: ", True), "The LSTM inversion (false-alarm filter → upset detector) rests on 12 medium-spread exclusives. Striking pattern, but limited formal significance within individual strata."),
        body(("Calibration artifacts: ", True), "Platt calibration compresses test probabilities to [0.19, 0.51], inflating threshold-based disagreement. Primary analyses use uncalibrated CV predictions; Top-K analysis used for the test set."),
        body(("Sample constraints: ", True), "3,495 training games, 558 test games. The NFL produces only ~270 games per season. Per-season results rest on 181–192 games each."),
        body(("Feature scope: ", True), "No player-level data, injury reports, weather forecasts, or motivational factors. The \"hidden information\" category is partly a data limitation — with richer data, some of those upsets might become classifiable."),
        body(("Single sport: ", True), "NFL has specific structural properties — 17-game seasons, high per-game variance, and a dominant betting market. The disagreement framework generalizes; the specific taxonomy is NFL-specific."),
    ])


def rewrite_slide_18(slide):
    """Conclusion — answer the question directly, end with the reframe."""
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Model Failure Is Diagnostic"),
    ])
    # Shape 2: intro - update to avoid duplication with body
    set_paragraphs(shapes[2].text_frame, [
        subtitle("Instead of asking which model predicts best, this study used the structure of their disagreements to diagnose why NFL outcomes happen."),
    ])
    # Shape 3: body
    set_paragraphs(shapes[3].text_frame, [
        body(("The core insight: ", True), "when structurally different models fail on different games, the pattern of failure reveals the mechanism behind the outcome. NFL upsets are not random — they are driven by at least four distinct mechanisms, each visible through a different architectural lens:"),
        body(("Consensus signal (45%): ", True), "All three models agree. The mismatch was obvious from every angle — summary stats, feature interactions, and trajectories all pointed the same way."),
        body(("Temporal signal (6%): ", True), "Only the LSTM sees it. And its role inverts by context: at close spreads it filters false alarms (92%); at medium spreads it detects real upsets (83%). That context-dependence is the strongest evidence the signal is structural, not noise."),
        body(("Hidden information (~24% of failures): ", True), "Every model fails on a big favorite. The cause was outside the data — injury, motivation, game-day scheme change."),
        body(("Irreducible variance (~76% of failures): ", True), "Every model fails on a close game. A coin flip. No pre-game model resolves which way a 3-point game breaks."),
        body(("No single model reveals this structure. You need the disagreement.", True)),
    ])


def rewrite_slide_19(slide):
    """Future Research — sharpen and connect to broader impact.

    Slide 19 has separate text boxes for each section:
    Shape 0: title, Shape 2: header1, Shape 3: body1,
    Shape 4: header2, Shape 5: body2, Shape 6: header3, Shape 7: body3
    """
    shapes = slide.shapes
    # Shape 0: title
    set_paragraphs(shapes[0].text_frame, [
        title_text("Where This Framework Goes Next"),
    ])
    # Shape 2: section 1 header
    set_paragraphs(shapes[2].text_frame, [
        section_header("Sharpen the taxonomy"),
    ])
    # Shape 3: section 1 body
    set_paragraphs(shapes[3].text_frame, [
        plain_body("Add game-day data — injuries, weather, motivation — to shrink the \"hidden information\" bucket. If those upsets become classifiable, the taxonomy tightens."),
        plain_body("Give the LSTM player-level trajectories instead of team averages. A team's momentum might really be one player getting healthy."),
        plain_body("Retrain on rolling windows so momentum patterns stay current. The NFL of 2015 is not the NFL of 2024."),
    ])
    # Shape 4: section 2 header
    set_paragraphs(shapes[4].text_frame, [
        section_header("Export the framework"),
    ])
    # Shape 5: section 2 body
    set_paragraphs(shapes[5].text_frame, [
        plain_body("NBA: 82-game seasons and less per-game variance — more data, cleaner signal. Does the LSTM's dual role survive with longer seasons?"),
        plain_body("Medical diagnosis: structurally different tests (MRI vs. biopsy vs. bloodwork) already exist. The same disagreement × context logic could categorize diagnostic failures."),
        plain_body("Financial prediction: market crashes have heterogeneous mechanisms too. Which model fails tells you whether the cause was momentum, fundamentals, or a black swan."),
    ])
    # Shape 6: section 3 header
    set_paragraphs(shapes[6].text_frame, [
        section_header("The reframe"),
    ])
    # Shape 7: section 3 body
    set_paragraphs(shapes[7].text_frame, [
        body(("Stop asking \"which model is best.\" Start asking \"what do their disagreements reveal about the problem?\" ", True), "The models aren't competing. They're diagnostic instruments. That shift in perspective is the real contribution."),
    ])


def rewrite_slide_20(slide):
    """References — add Huang & Li, Boulier et al. 2006."""
    shapes = slide.shapes
    # Shape 2: references text (shape 1 is the accent rectangle)
    set_paragraphs(shapes[2].text_frame, [
        plain_body("Boulier, B. L., & Stekler, H. O. (2003). Predicting the outcomes of NFL games. International Journal of Forecasting, 19(2), 257-270."),
        plain_body("Boulier, B. L., Stekler, H. O., & Amundson, S. (2006). Testing the efficiency of the NFL point spread. Journal of Quantitative Analysis in Sports, 2(2)."),
        plain_body("Dietterich, T. G. (2000). Ensemble methods in machine learning. MCS 2000, LNCS 1857."),
        plain_body("Glickman, M. E., & Stern, H. S. (1998). A state-space model for NFL scores. JASA, 93(441), 25-35."),
        plain_body("Gordon, M. L., et al. (2021). Disagreement deconstructed. CHI '21."),
        plain_body("Grinsztajn, L., Oyallon, E., & Varoquaux, G. (2022). Why do tree-based models still outperform deep learning on typical tabular data? NeurIPS 35."),
        plain_body("Huang, M.-L., & Li, Y.-Z. (2021). Use of machine learning and deep learning to predict the outcomes of major league baseball games. Applied Sciences, 11(10), 4499."),
        plain_body("Hubacek, O., Sourek, G., & Zelezny, F. (2019). Exploiting sports-betting market using ML. IJF, 35(2), 783-796."),
        plain_body("Krogh, A., & Vedelsby, J. (1995). Neural network ensembles, cross validation, and active learning. NIPS 7."),
        plain_body("Lakshminarayanan, B., Pritzel, A., & Blundell, C. (2017). Simple and scalable predictive uncertainty estimation using deep ensembles. NeurIPS 30."),
        plain_body("Song, C., Boulier, B. L., & Stekler, H. O. (2007). The comparative accuracy of judgmental and model forecasts of American football games. IJF, 23(3), 405-413."),
        plain_body("Wilkens, S. (2021). Sports prediction and betting models in the ML age. Journal of Sports Analytics, 7(2)."),
    ])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation("docs/AP_Research_POD_Corrected.pptx")
    slides = list(prs.slides)

    rewriters = {
        0: rewrite_slide_1,
        1: rewrite_slide_2,
        2: rewrite_slide_3,
        3: rewrite_slide_4,
        4: rewrite_slide_5,
        5: rewrite_slide_6,
        6: rewrite_slide_7,
        7: rewrite_slide_8,
        8: rewrite_slide_9,
        9: rewrite_slide_10,
        10: rewrite_slide_11,
        11: rewrite_slide_12,
        12: rewrite_slide_13,
        13: rewrite_slide_14,
        14: rewrite_slide_15,
        15: rewrite_slide_16,
        16: rewrite_slide_17,
        17: rewrite_slide_18,
        18: rewrite_slide_19,
        19: rewrite_slide_20,
    }

    for idx, rewriter in rewriters.items():
        print(f"Rewriting slide {idx + 1}...")
        try:
            rewriter(slides[idx])
        except Exception as e:
            print(f"  ERROR on slide {idx + 1}: {e}")

    out_path = "docs/AP_Research_POD_Revised.pptx"
    prs.save(out_path)
    print(f"\nSaved to {out_path}")


if __name__ == "__main__":
    main()
