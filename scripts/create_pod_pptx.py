"""Create AP Research POD presentation as .pptx matching the uploaded style."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Style constants (matching uploaded deck) ──────────────────────────
TEAL = RGBColor(0x4E, 0x9E, 0xAD)       # Teal accent / subheadings
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)    # Main header black
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)    # Body text gray
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)  # Subtitle / caption gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER_BG = RGBColor(0x4E, 0x9E, 0xAD)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF2)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


def add_accent_line(slide, left=Inches(0.8), top=None, width=Inches(2.2)):
    """Add the teal accent line under the header."""
    if top is None:
        top = Inches(1.15)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    return shape


def set_cell_text(cell, text, font_size=14, bold=False, color=DARK_TEXT,
                  alignment=PP_ALIGN.LEFT):
    """Set text in a table cell with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_BODY
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, rows_data, left=Inches(0.8), top=Inches(2.2),
              col_widths=None, font_size=13):
    """Add a styled table to a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    width = sum(col_widths) if col_widths else Inches(11)
    height = Inches(0.4 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            is_header = r_idx == 0
            color = WHITE if is_header else DARK_TEXT
            bold = is_header
            set_cell_text(cell, cell_text, font_size=font_size,
                          bold=bold, color=color)
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.background()
    return table_shape


def add_header(slide, text, top=Inches(0.4), left=Inches(0.8)):
    """Add a bold header to a slide."""
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT_TITLE
    return txBox


def add_subheader(slide, text, top=Inches(1.4), left=Inches(0.8)):
    """Add a teal subheader."""
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = TEAL
    run.font.name = FONT_BODY
    return txBox


def add_body_text(slide, text, top=Inches(2.0), left=Inches(0.8),
                  width=Inches(11.5), font_size=16):
    """Add gray body text, supporting bullet points via '• ' prefix."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        stripped = line.strip()
        if not stripped:
            p.space_after = Pt(6)
            continue

        # Detect bold segments: **text**
        if stripped.startswith("• "):
            stripped = stripped[2:]
            p.level = 0
            p.space_before = Pt(4)

        parts = _parse_bold(stripped)
        for text_part, is_bold in parts:
            run = p.add_run()
            run.text = text_part
            run.font.size = Pt(font_size)
            run.font.color.rgb = GRAY_TEXT
            run.font.name = FONT_BODY
            run.font.bold = is_bold
        p.space_after = Pt(4)
    return txBox


def _parse_bold(text):
    """Parse **bold** segments in text. Returns list of (text, is_bold)."""
    parts = []
    while "**" in text:
        idx = text.index("**")
        if idx > 0:
            parts.append((text[:idx], False))
        text = text[idx + 2:]
        end_idx = text.find("**")
        if end_idx == -1:
            parts.append((text, True))
            return parts
        parts.append((text[:end_idx], True))
        text = text[end_idx + 2:]
    if text:
        parts.append((text, False))
    return parts if parts else [(text, False)]


def add_section_label(slide, text, top=Inches(2.0), left=Inches(0.8)):
    """Add a teal bold section label (e.g., 'Sports Prediction')."""
    txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = FONT_BODY
    return txBox


# ══════════════════════════════════════════════════════════════════════
#  BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]  # blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The Anatomy\nof NFL Upsets"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT_TITLE

    add_accent_line(slide, top=Inches(4.0), width=Inches(2.5))

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Comparing Logistic Regression, XGBoost, and LSTM"
    run2.font.size = Pt(20)
    run2.font.color.rgb = TEAL
    run2.font.name = FONT_BODY

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Wil Fowler  |  AP Research  |  2025-2026"
    run3.font.size = Pt(16)
    run3.font.color.rgb = LIGHT_GRAY
    run3.font.name = FONT_BODY

    # ── Slide 2: Key Terms ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Key Terms You Need to Know")
    add_accent_line(slide)
    add_subheader(slide, "Here are the key terms that come up throughout this presentation.")
    add_body_text(slide, (
        "• **AUC (Area Under the Curve):** A score from 0 to 1 measuring how well a model ranks outcomes. "
        "0.50 = random guessing, 1.0 = perfect. Think of it like a grade.\n"
        "• **p-value:** The probability a result happened by pure luck. Below .05 = likely real. "
        "Above .05 = cannot rule out chance.\n"
        "• **Bootstrap Confidence Interval:** A range computed by resampling the data thousands of times. "
        "If the interval contains zero, the difference is not statistically significant.\n"
        "• **Cross-Validation:** Training on older data, testing on newer data the model has never seen. "
        "Prevents memorizing answers.\n"
        "• **LSTM:** A model that processes games in sequence (week 1, week 2, week 3...) instead of "
        "looking at each game alone. Detects momentum and fatigue.\n"
        "• **Point Spread:** How many points oddsmakers think the favorite will win by. A spread of "
        "7 = the market expects a touchdown margin.\n"
        "• **Spread Ablation:** Removing the point spread from the models to see how much they depend "
        "on it versus the raw stats."
    ), top=Inches(1.9), font_size=15)

    # ── Slide 3: Research Question ────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Research Question")
    add_accent_line(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("To what extent can disagreement patterns among logistic regression, "
                "XGBoost, and LSTM models, combined with spread ablation analysis, "
                "reveal distinct categories of NFL upsets?")
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.color.rgb = TEAL
    run.font.name = FONT_BODY

    add_body_text(slide, (
        "• NFL upsets happen when a 3+ point underdog wins. That's about 29% of games.\n"
        "• Most models try to predict wins and losses. This study looks at **why** certain upsets "
        "happen by analyzing which models get them right and which don't.\n"
        "• Spread ablation removes betting line data to test whether models actually learn "
        "something on their own or just echo the market."
    ), top=Inches(3.2), font_size=16)

    # ── Slide 4: Literature Review ────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Literature Review")
    add_accent_line(slide)
    add_section_label(slide, "Sports Prediction", top=Inches(1.5))
    add_body_text(slide, (
        "• Logistic regression and statistical models remain competitive with complex ML for "
        "NFL game prediction (Boulier & Stekler, 2003)\n"
        "• Tree-based models outperform deep learning on typical tabular data but not necessarily "
        "on sequential data (Grinsztajn et al., 2022)\n"
        "• Deep learning including LSTMs shows promise for capturing sequential patterns in sports "
        "outcomes (Huang & Li, 2021)"
    ), top=Inches(1.9), font_size=15)
    add_section_label(slide, "Market Efficiency", top=Inches(3.5))
    add_body_text(slide, (
        "• NFL betting lines outperform both expert judgment and statistical models "
        "(Song et al., 2007; Wilkens, 2021)\n"
        "• Decorrelating predictions from bookmaker odds is necessary to find genuine signal "
        "(Hubacek et al., 2019)\n"
        "• Market inefficiencies exist in specific contexts but are generally small and hard to "
        "exploit (Boulier et al., 2006)"
    ), top=Inches(3.9), font_size=15)
    add_section_label(slide, "Model Disagreement", top=Inches(5.5))
    add_body_text(slide, (
        "• Per-instance disagreement between models is an unbiased estimate of error variance "
        "(Gordon et al., 2021)\n"
        "• Architecturally diverse models pick up on different data patterns (Dietterich, 2000)\n"
        "• **No prior work uses model disagreement to categorize upsets into distinct types**"
    ), top=Inches(5.9), font_size=15)

    # ── Slide 5: Gap ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Gap in the Literature")
    add_accent_line(slide)
    add_subheader(slide, ("No one has used model disagreement patterns to classify NFL upsets "
                          "into distinct types. Existing research treats all upsets the same."))
    add_body_text(slide, (
        "• Prior studies focus on binary win/loss prediction, not understanding why "
        "specific upsets occur\n"
        "• Ensemble disagreement research tries to improve accuracy, not use disagreement "
        "as a diagnostic tool\n"
        "• No one has tested whether prediction models learn independent signals or just "
        "reflect what the betting market already knows"
    ), top=Inches(3.0), font_size=16)

    # ── Slide 6: Method Overview ──────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Method Overview")
    add_accent_line(slide)
    add_subheader(slide, "Quantitative, quasi-experimental with observational data and ablation analysis")
    add_body_text(slide, (
        "1. Collected 18 NFL regular seasons (2005-2022). Held out 2023-2025 as a true test set.\n\n"
        "2. Gave each model the **same data in a different form** — like three analysts watching\n"
        "     the same games but reading different scouting reports:\n"
        "     LR reads **\"The Summary\"** (46 stats). XGB reads **\"The Details\"** (70 features with\n"
        "     game-by-game lags). LSTM watches **\"The Movie\"** (raw 8-game sequences).\n\n"
        "3. Trained all three on identical folds (6-fold expanding-window CV). Same data, same\n"
        "     splits — so when they disagree, it's the architecture talking, not the data.\n\n"
        "4. Classified every game by **who got it right and who didn't**. The pattern of agreement\n"
        "     and disagreement is the diagnostic tool.\n\n"
        "5. Removed the point spread and reran everything. If a model collapses, it was just\n"
        "     echoing the market. If it survives, it learned something real."
    ), top=Inches(2.0), font_size=16)

    # ── Slide 7: Data & Features ──────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Data & Multi-Representation Design")
    add_accent_line(slide)

    # Dataset box (left)
    add_section_label(slide, "Dataset", top=Inches(1.5), left=Inches(0.8))
    add_body_text(slide, (
        "Training: 3,495 games (2005-2022)\n"
        "Test: 558 games (2023-2025)\n"
        "Upset-eligible: spread of 3+ points\n"
        "Base rate: ~30% upsets"
    ), top=Inches(2.0), left=Inches(0.8), width=Inches(4.5), font_size=15)

    # Representation table
    add_table(slide, [
        ["Model", "Representation", "Features", "What It Sees"],
        ["LR", '"The Summary"', "46", "Rolling averages, differentials, market, Elo"],
        ["XGBoost", '"The Details"', "70", "LR's 46 + 24 per-game lag stats (last 3 games)"],
        ["LSTM", '"The Movie"', "14×8 + 10", "Raw game-by-game sequences + matchup context"],
    ], top=Inches(3.8), left=Inches(0.8),
       col_widths=[Inches(1.5), Inches(2.2), Inches(1.2), Inches(6.6)],
       font_size=13)

    # Models section
    add_section_label(slide, "Why these three?", top=Inches(5.6), left=Inches(0.8))
    add_body_text(slide, (
        "**LR** sees linear relationships — is the spread just wrong? "
        "**XGBoost** sees interactions — does a bad recent game + short week + road = danger? "
        "**LSTM** sees trajectories — is this team on the way up or falling apart?"
    ), top=Inches(6.0), left=Inches(0.8), width=Inches(11.5), font_size=15)

    # ── Slide 8: CV Performance ───────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Results: Cross-Validation Performance")
    add_accent_line(slide)
    add_table(slide, [
        ["Model", "AUC", "Brier Score", "Log Loss"],
        ["Logistic Regression", ".650", ".197", ".581"],
        ["LSTM", ".641", ".199", ".583"],
        ["XGBoost", ".638", ".199", ".586"],
    ], top=Inches(1.6), left=Inches(1.5),
       col_widths=[Inches(3), Inches(2.2), Inches(2.2), Inches(2.2)],
       font_size=14)

    # Caption
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "6-fold expanding-window CV, validation years 2017-2022, N = 1,162 games"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = FONT_BODY

    add_body_text(slide, (
        "• All three models are **statistically indistinguishable**. Bootstrap 95% CIs on all "
        "pairwise AUC differences contain zero.\n"
        "• This validates the disagreement framework: disagreement reflects **different "
        "perspectives**, not differences in model quality.\n"
        "• LR and XGBoost predictions correlated at .874. The LSTM correlated moderately "
        "with both (.784 and .699), confirming it captures partially distinct information.\n"
        "• The XGB-LSTM correlation (.699) is the lowest, reflecting the greatest structural "
        "distance between interaction-based and sequence-based processing."
    ), top=Inches(3.9), font_size=15)

    # ── Slide 9: Disagreement Categories ──────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Results: Disagreement Categories")
    add_accent_line(slide)
    add_table(slide, [
        ["Category", "N", "% of Games", "Upset Rate", "Meaning"],
        ["All correct", "528", "45.4%", "36.7%", "Strong cross-architectural signal"],
        ["All wrong", "340", "29.3%", "20.9%", "Outside model capabilities"],
        ["Only LSTM", "65", "5.6%", "18.5%", "Temporal signal (unique to LSTM)"],
        ["LR + XGB only", "78", "6.7%", "29.5%", "Static models agree, LSTM disagrees"],
        ["Only XGB", "48", "4.1%", "25.0%", "Non-linear interaction pattern"],
        ["LR + LSTM", "48", "4.1%", "33.3%", "Linear + temporal agreement"],
        ["Only LR", "28", "2.4%", "28.6%", "Linear spread mispricing"],
        ["XGB + LSTM", "27", "2.3%", "33.3%", "Non-linear + temporal"],
    ], top=Inches(1.5), left=Inches(0.5),
       col_widths=[Inches(2.2), Inches(0.9), Inches(1.6), Inches(1.6), Inches(5.2)],
       font_size=12)

    # Caption
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(11), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CV set: 1,162 games, threshold = base upset rate (~0.30)"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = FONT_BODY

    add_body_text(slide, (
        "• 74.7% of games fall into all-agree categories (all correct + all wrong).\n"
        "• **Key finding:** Of 65 LSTM-exclusive games, 53 (81.5%) are **non-upset rejections** — "
        "the LSTM moderates false alarms from the other models. Only 12 are upsets caught.\n"
        "• This bias is statistically significant (one-sided binomial p = 0.029)."
    ), top=Inches(5.4), font_size=15)

    # ── Slide 10: Spread-Stratified (NEW) ─────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Results: The LSTM's Role Depends on Spread")
    add_accent_line(slide)
    add_subheader(slide, "This is the study's most diagnostic finding.")
    add_table(slide, [
        ["Spread Bucket", "LSTM Exclusives", "Upsets Caught", "Non-Upsets Rejected", "Primary Role"],
        ["Small (3-6.5 pts)", "60", "5 (8%)", "55 (92%)", "False-alarm filter"],
        ["Medium (7-13.5 pts)", "12", "10 (83%)", "2 (17%)", "Upset detector"],
        ["Large (14+ pts)", "0", "—", "—", "No exclusive value"],
    ], top=Inches(2.2), left=Inches(0.8),
       col_widths=[Inches(2.5), Inches(2.2), Inches(2), Inches(2.5), Inches(2.3)],
       font_size=13)

    add_body_text(slide, (
        "• **Close games (3-6.5 pts):** LR and XGB look at the stats and say \"upset.\" "
        "The LSTM watches the last 8 games and says \"no — this team is trending the wrong way.\" "
        "It's right **92%** of the time. The movie tells a different story than the snapshot.\n"
        "• **Bigger underdogs (7-13.5 pts):** Now the LSTM flips. It sees a team getting "
        "hotter week by week — improving EPA, winning close games, building momentum. The "
        "summary stats say \"no chance.\" The trajectory says **\"watch this team.\"** "
        "It catches 83% of these as real upsets.\n"
        "• This inversion is the study's key finding. The LSTM isn't just \"better\" or "
        "\"worse\" — it reads the game differently depending on context."
    ), top=Inches(4.2), font_size=15)

    # ── Slide 11: Spread Ablation ─────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Results: Spread Ablation")
    add_accent_line(slide)
    add_table(slide, [
        ["Model", "AUC With Spread", "AUC Without", "Delta", "p"],
        ["LR", ".650", ".571", "-.079", "< .001"],
        ["XGBoost", ".638", ".566", "-.072", "< .001"],
        ["LSTM", ".641", ".574", "-.067", "< .001"],
    ], top=Inches(1.6), left=Inches(1.5),
       col_widths=[Inches(2), Inches(2.2), Inches(2.2), Inches(1.5), Inches(1.5)],
       font_size=14)

    add_body_text(slide, (
        "• All three models depend heavily on spread data. Bootstrap CIs confirm all deltas "
        "are statistically significant.\n"
        "• The LSTM degrades least (-.067) and **becomes the strongest model** without spread "
        "(.574 vs .571 and .566). The ranking reverses.\n"
        "• But the LSTM's smaller delta is **not statistically significant** vs LR's delta "
        "(bootstrap CI contains zero). Suggestive, not proven.\n"
        "• Without spread, LSTM-exclusive predictions **double** (5.6% → 11.0%) and upsets "
        "caught increase from 12 to 33. The spread was masking temporal signal.\n"
        "• LR-XGB correlation dropped from .874 to .742. The spread was the shared anchor "
        "holding model predictions together."
    ), top=Inches(3.5), font_size=15)

    # ── Slide 12: Test Set ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Results: Out-of-Sample Test (2023-2025)")
    add_accent_line(slide)

    # CV-to-test table
    add_table(slide, [
        ["Model", "CV AUC", "Test AUC", "Gap"],
        ["XGBoost", ".638", ".576", "-.062"],
        ["LR", ".650", ".562", "-.088"],
        ["LSTM", ".641", ".524", "-.117"],
    ], top=Inches(1.6), left=Inches(0.8),
       col_widths=[Inches(2), Inches(1.5), Inches(1.5), Inches(1.5)],
       font_size=14)

    # Top-K table
    txBox = slide.shapes.add_textbox(Inches(7), Inches(1.35), Inches(5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Top-K Hit Rates (base rate = 28.5%)"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = TEAL
    run.font.name = FONT_BODY

    add_table(slide, [
        ["K", "LR", "XGB", "LSTM", "Ensemble"],
        ["10", "50%", "60%", "40%", "50%"],
        ["20", "40%", "45%", "35%", "45%"],
        ["50", "38%", "44%", "28%", "34%"],
    ], top=Inches(1.6), left=Inches(7),
       col_widths=[Inches(0.8), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.4)],
       font_size=13)

    add_body_text(slide, (
        "• The LSTM shows the **largest** generalization gap (-.117). Temporal patterns "
        "from 2005-2022 do not fully transfer to 2023-2025.\n"
        "• XGBoost generalizes best (-.062). Non-linear interaction patterns are the most "
        "temporally stable signal.\n"
        "• LSTM correlation with static models drops dramatically on test: "
        "LR .784 → .429, XGB .699 → .408. It diverges more in truly out-of-sample data.\n"
        "• XGBoost's top 10 predictions hit at **60% (2.1x lift)**. Signal concentrates in "
        "highest-confidence predictions and fades quickly."
    ), top=Inches(3.8), font_size=15)

    # ── Slide 13: Four Types ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Not All Upsets Are Created Equal")
    add_accent_line(slide)
    add_subheader(slide,
                  "If upsets were random, every model would fail the same way. They don't. "
                  "The pattern of who gets it right tells you why it happened.")
    add_body_text(slide, (
        "• **All three agree → The mismatch was obvious from every angle.** The Summary, "
        "The Details, and The Movie all told the same story. A team that's better on paper, "
        "in the interactions, and on the trajectory. 45% of games. These are the structurally "
        "readable outcomes — the games oddsmakers price well.\n\n"
        "• **Only the LSTM sees it → The movie told a story the snapshot couldn't.** "
        "A defense giving up more yards each week. An offense that just found its rhythm. "
        "A team coming off a bye after three straight losses. The summary says \"average.\" "
        "The trajectory says \"momentum.\" 6% of games — small, but structurally distinct.\n\n"
        "• **All three fail at big spreads → Something happened we couldn't see.** "
        "A QB injury at warmups. A team resting starters with a playoff bye locked up. A "
        "locker room that imploded Tuesday. No historical stat captures game-day surprises.\n\n"
        "• **All three fail at small spreads → A coin flip that landed tails.** "
        "3-point games are toss-ups. A tipped pass, a missed field goal, a bad spot on "
        "4th-and-inches. Every model knew it was close. No one can predict which way close "
        "games break."
    ), top=Inches(2.2), font_size=15)

    # ── Slide 14: LSTM Paradox ────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Momentum Is Real — but It Doesn't Stay the Same")
    add_accent_line(slide)
    add_subheader(slide,
                  "The LSTM matched the other models in training but showed the largest "
                  "drop on truly future data. That's not a failure — it's a finding.")

    add_section_label(slide, "In the training era (2005-2022):", top=Inches(2.1))
    add_body_text(slide, (
        "• LSTM AUC .641 — statistically tied with LR (.650) and XGB (.638)\n"
        "• The temporal signal is real: it adds value to the ensemble and its exclusive "
        "catches double when the spread is removed\n"
        "• Momentum patterns from this era — hot streaks, schedule fatigue, form "
        "trajectories — are genuinely predictive"
    ), top=Inches(2.5), font_size=15)

    add_section_label(slide, "In the future (2023-2025):", top=Inches(4.1))
    add_body_text(slide, (
        "• LSTM drops to .524 — largest gap of any model (-.117)\n"
        "• It trails LR and XGB in all three test seasons. Not one weird year — persistent.\n"
        "• The correlation between LSTM and the static models collapses: .78 → .43"
    ), top=Inches(4.5), font_size=15)

    add_section_label(slide, "What this means about the NFL:", top=Inches(5.7))
    add_body_text(slide, (
        "Momentum patterns are **era-specific**. How a hot streak translated into upset "
        "probability in 2015 doesn't translate the same way in 2024. Coaching changes, "
        "rule changes, the way offenses evolve — the temporal dynamics of the NFL shift "
        "faster than the static relationships. That's a finding about the sport."
    ), top=Inches(6.1), font_size=15)

    # ── Slide 15: Spread Ablation Interpretation ──────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "What the Spread Ablation Tells Us")
    add_accent_line(slide)
    add_subheader(slide,
                  "The spread is the answer sheet. When I took it away, I found out who was "
                  "actually learning and who was just copying.")
    add_body_text(slide, (
        "• LR lost .079 AUC, XGBoost lost .072, LSTM lost .067. All statistically significant.\n"
        "• LR's biggest coefficient? The spread itself — 5x larger than any team stat. "
        "The \"model\" was mostly just reading the market's homework.\n"
        "• **But the LSTM kept the most signal.** Without the spread, the ranking flips: "
        "LSTM (.574) > LR (.571) > XGB (.566). The Movie captures something real that the "
        "market encodes differently than raw trajectory.\n"
        "• Model agreement collapsed: LR-XGB correlation .874 → .742, XGB-LSTM .699 → .419. "
        "The spread was the shared anchor. Without it, the architectures scatter — revealing "
        "how differently they actually read the game.\n"
        "• LSTM-exclusive catches **doubled** without spread (12 → 33 upsets). The betting "
        "line was masking genuine temporal signal.\n"
        "• **Bottom line:** The market is good. Oddsmakers have already digested everything "
        "in publicly available stats. But the LSTM finds something the line encodes "
        "differently — trajectory, momentum, the shape of recent results."
    ), top=Inches(2.3), font_size=15)

    # ── Slide 16: Bottom Line ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Why Disagreement Proves Structure")
    add_accent_line(slide)
    add_subheader(slide,
                  "Three models. Same data. Different architectures. If upsets were random, "
                  "they'd fail randomly. They don't.")
    add_body_text(slide, (
        "• **The Summary and The Details agree on 87% of predictions.** When the static "
        "picture is clear — one team is better on paper — both linear and non-linear "
        "models see it. When they're both wrong, the upset came from outside the numbers.\n\n"
        "• **The Movie tells a different story — and it depends on context.** "
        "In close games, the LSTM corrects the other two: \"the stats say upset, but watch "
        "the trajectory — this team is fading.\" In bigger mismatches, it catches what they "
        "miss: \"the numbers say no, but this team has been building for three weeks.\"\n\n"
        "• **Remove the betting line and the masks come off.** The spread was the anchor — "
        "the thing all three models leaned on. Without it, agreement collapses, the LSTM's "
        "unique signal doubles, and the three architectures reveal how differently they "
        "actually read the game.\n\n"
        "• **Momentum is real, but it changes shape.** The LSTM finds temporal patterns that "
        "improve predictions during training. But the NFL evolves — coaching, rules, "
        "offensive philosophy — and those patterns don't stay the same. That tells us "
        "something about the sport itself: **momentum matters, but its signature shifts.**"
    ), top=Inches(2.2), font_size=15)

    # ── Slide 17: Limitations ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Limitations")
    add_accent_line(slide)
    add_body_text(slide, (
        "• **Small cell sizes in spread strata:** The LSTM inversion (false-alarm filter → upset "
        "detector) rests on 12 medium-spread exclusives. Striking pattern, but limited formal "
        "significance within strata.\n\n"
        "• **Calibration artifacts:** Platt calibration compresses test probabilities to [0.19, 0.51], "
        "inflating threshold-based disagreement. Primary analyses use uncalibrated CV predictions; "
        "rank-based (Top-K) analysis used for test set.\n\n"
        "• **Sample size:** 3,495 training games, 558 test games. NFL produces ~270 games/season. "
        "Per-season results based on 181-192 games each.\n\n"
        "• **Feature scope:** No player-level data, injury reports, weather forecasts, or motivational "
        "factors. The \"hidden information\" category is partly a data availability limitation.\n\n"
        "• **Single sport:** NFL has specific structural properties (17-game seasons, high variance, "
        "dominant betting market). The disagreement framework generalizes; the specific taxonomy "
        "is NFL-specific."
    ), top=Inches(1.6), font_size=15)

    # ── Slide 18: Conclusion ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Conclusion: Model Failure Is Diagnostic")
    add_accent_line(slide)
    add_subheader(slide,
                  "Instead of asking which model predicts best, I used the structure of "
                  "their disagreements to diagnose why NFL outcomes happen.")

    add_body_text(slide, (
        "The big idea: **when structurally different models fail on different games, "
        "the pattern of failure reveals the mechanism behind the outcome.**\n\n"
        "That's what this study found. NFL upsets are not random — they are driven by "
        "at least four distinct mechanisms, each visible through a different lens:\n\n"
        "• **Static mismatch** — all three models agree. The game was readable from any "
        "angle: summary stats, feature interactions, and trajectories all pointed the same "
        "way. The market prices these well. (~45% of games)\n\n"
        "• **Temporal momentum** — only the LSTM sees it. Something was changing "
        "week-to-week that snapshots missed. And the LSTM's role **inverts by context**: "
        "at close spreads it filters false alarms (92%); at bigger spreads it detects "
        "real upsets (83%). That context-dependence proves the signal is structural, "
        "not noise. (~6% of games)\n\n"
        "• **Hidden information** — every model fails on a big favorite. The cause was "
        "a game-day injury, motivation, or scheme change — outside the data. (~24% of "
        "all_wrong at medium/large spreads)\n\n"
        "• **Irreducible variance** — every model fails on a close game. A coin flip. "
        "(~76% of all_wrong at small spreads)\n\n"
        "No single model reveals this structure. You need the disagreement."
    ), top=Inches(1.9), font_size=14)

    # ── Slide 19: Future Research & Implications ────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "What Comes Next")
    add_accent_line(slide)

    add_section_label(slide, "Make the taxonomy sharper", top=Inches(1.5))
    add_body_text(slide, (
        "• Add **game-day data** — injuries, weather, motivation — to shrink the \"hidden "
        "information\" bucket. If those upsets become explainable, the taxonomy tightens.\n"
        "• Give the LSTM **player-level trajectories** instead of team averages. A team's "
        "momentum might really be one player getting healthy — the LSTM should see that.\n"
        "• **Retrain on rolling windows** so the momentum patterns stay current. The NFL "
        "of 2015 is not the NFL of 2024."
    ), top=Inches(1.9), font_size=15)

    add_section_label(slide, "Export the framework", top=Inches(3.7))
    add_body_text(slide, (
        "• **NBA** has 82-game seasons and less per-game variance — more data, cleaner signal. "
        "Does the LSTM's temporal role survive with more games per season?\n"
        "• **Medical diagnosis** already uses structurally different tests (MRI vs. biopsy vs. "
        "bloodwork). The same disagreement × context logic could categorize diagnostic failures.\n"
        "• **Financial prediction** — market crashes have different mechanisms too. "
        "Which model fails tells you whether it was momentum, fundamentals, or a black swan."
    ), top=Inches(4.1), font_size=15)

    add_section_label(slide, "The reframe", top=Inches(5.9))
    add_body_text(slide, (
        "Stop asking **\"which model is best.\"** Start asking **\"what do their "
        "disagreements tell me about the problem?\"** The models aren't competing. "
        "They're diagnostic instruments."
    ), top=Inches(6.25), font_size=16)

    # ── Slide 20: References ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_header(slide, "References")
    add_accent_line(slide)
    add_body_text(slide, (
        "Boulier, B. L., & Stekler, H. O. (2003). Predicting the outcomes of NFL games. "
        "International Journal of Forecasting, 19(2), 257-270.\n\n"
        "Dietterich, T. G. (2000). Ensemble methods in machine learning. MCS 2000, LNCS 1857.\n\n"
        "Glickman, M. E., & Stern, H. S. (1998). A state-space model for NFL scores. "
        "JASA, 93(441), 25-35.\n\n"
        "Gordon, M. L., et al. (2021). Disagreement deconstructed. CHI '21.\n\n"
        "Grinsztajn, L., Oyallon, E., & Varoquaux, G. (2022). Why do tree-based models still "
        "outperform deep learning on typical tabular data? NeurIPS 35.\n\n"
        "Hubacek, O., Sourek, G., & Zelezny, F. (2019). Exploiting sports-betting market "
        "using ML. IJF, 35(2), 783-796.\n\n"
        "Krogh, A., & Vedelsby, J. (1995). Neural network ensembles, cross validation, "
        "and active learning. NIPS 7.\n\n"
        "Lakshminarayanan, B., Pritzel, A., & Blundell, C. (2017). Simple and scalable "
        "predictive uncertainty estimation using deep ensembles. NeurIPS 30.\n\n"
        "Wilkens, S. (2021). Sports prediction and betting models in the ML age. "
        "J Sports Analytics, 7(2)."
    ), top=Inches(1.5), font_size=13)

    # ── Save ──────────────────────────────────────────────────────────
    out_path = "docs/AP_Research_POD_Corrected.pptx"
    prs.save(out_path)
    print(f"Saved to {out_path}")
    return out_path


if __name__ == "__main__":
    build_presentation()
